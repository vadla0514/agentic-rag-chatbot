# ingestion_agent.py

import os
import pandas as pd
import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation

# Extract text from various file types
def extract_text_from_file(file_path):
    ext = os.path.splitext(file_path)[-1].lower()

    if ext == '.pdf':
        return extract_text_from_pdf(file_path)
    elif ext == '.docx':
        return extract_text_from_docx(file_path)
    elif ext in ['.txt', '.md']:
        return extract_text_from_txt(file_path)
    elif ext == '.csv':
        return extract_text_from_csv(file_path)
    elif ext == '.pptx':
        return extract_text_from_pptx(file_path)
    else:
        return ""

# File-type specific handlers
def extract_text_from_pdf(path):
    text = ""
    doc = fitz.open(path)
    for page in doc:
        text += page.get_text()
    return text

def extract_text_from_docx(path):
    doc = Document(path)
    return "\n".join([para.text for para in doc.paragraphs])

def extract_text_from_txt(path):
    with open(path, 'r', encoding='utf-8') as f:
        return f.read()

def extract_text_from_csv(path):
    df = pd.read_csv(path)
    return df.to_string(index=False)

def extract_text_from_pptx(path):
    prs = Presentation(path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

# MCP message constructor
def create_mcp_message(sender, receiver, trace_id, all_text):
    return {
        "sender": sender,
        "receiver": receiver,
        "type": "CONTEXT_INIT",
        "trace_id": trace_id,
        "payload": {
            "raw_text": all_text
        }
    }

# Agent logic
def ingest_files(file_paths, trace_id="trace-001"):
    all_text = ""
    for file_path in file_paths:
        content = extract_text_from_file(file_path)
        all_text += content + "\n\n"

    return create_mcp_message(
        sender="IngestionAgent",
        receiver="RetrievalAgent",
        trace_id=trace_id,
        all_text=all_text
    )
